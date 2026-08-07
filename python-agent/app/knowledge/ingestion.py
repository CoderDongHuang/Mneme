import csv
import hashlib
import zipfile
from collections import Counter
from html.parser import HTMLParser
from pathlib import Path

from langchain_core.documents import Document

from app.core.config import settings
from app.core.logging import setup_logger
from app.knowledge.chunking import chunk_documents
from app.knowledge.vector_store import vector_store


logger = setup_logger("ingestion")
SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".xlsm",
    ".csv",
    ".md",
    ".markdown",
    ".txt",
    ".html",
    ".htm",
}


class _TextHTMLParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data: str) -> None:
        text = data.strip()
        if text:
            self.parts.append(text)


def _document(content: str, source: str, **metadata: object) -> Document:
    return Document(page_content=content, metadata={"source": source, **metadata})


def _parse_pdf(path: Path, source: str) -> list[Document]:
    import fitz

    pdf = fitz.open(str(path))
    documents: list[Document] = []
    page_lines: list[list[str]] = []
    page_images: list[tuple[int, bytes]] = []

    for page_index, page in enumerate(pdf, start=1):
        blocks = sorted(
            page.get_text("blocks"),
            key=lambda block: (round(float(block[1]) / 12), float(block[0])),
        )
        lines = [
            line.strip()
            for block in blocks
            for line in str(block[4]).splitlines()
            if line.strip()
        ]
        text = "\n".join(lines).strip()

        if len(text) < settings.pdf_min_text_chars and settings.ocr_enabled:
            import pytesseract
            from PIL import Image

            pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            image = Image.frombytes(
                "RGB", [pixmap.width, pixmap.height], pixmap.samples
            )
            text = pytesseract.image_to_string(
                image, lang=settings.ocr_languages
            ).strip()
            lines = [line.strip() for line in text.splitlines() if line.strip()]
        page_lines.append(lines)

        if (
            settings.multimodal_enabled
            and len(page_images) < settings.multimodal_max_images
            and (page.get_images(full=True) or page.get_drawings())
        ):
            rendered = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
            page_images.append((page_index, rendered.tobytes("png")))

    edge_lines: Counter[str] = Counter()
    for lines in page_lines:
        for line in lines[:2] + lines[-2:]:
            if 2 <= len(line) <= 160:
                edge_lines[line] += 1
    repeated = {
        line
        for line, count in edge_lines.items()
        if len(page_lines) >= 3 and count >= max(3, len(page_lines) // 2)
    }

    for page_index, lines in enumerate(page_lines, start=1):
        text = "\n".join(line for line in lines if line not in repeated).strip()
        if text:
            documents.append(
                _document(
                    text,
                    source,
                    page=page_index,
                    chunk_type="text",
                    parser="pymupdf_layout",
                )
            )

    try:
        import pdfplumber

        with pdfplumber.open(str(path)) as pdf:
            for page_index, page in enumerate(pdf.pages, start=1):
                for table_index, table in enumerate(
                    page.extract_tables() or [], start=1
                ):
                    rows = [
                        " | ".join("" if cell is None else str(cell) for cell in row)
                        for row in table
                    ]
                    content = "\n".join(rows).strip()
                    if content:
                        documents.append(
                            _document(
                                content,
                                source,
                                page=page_index,
                                chunk_type="table",
                                section=f"第 {page_index} 页表格 {table_index}",
                            )
                        )
    except ImportError:
        logger.info("未安装 pdfplumber，PDF 表格将按普通文本处理")
    except Exception as error:
        logger.warning("PDF 表格解析失败，继续使用文本结果: %s", error)

    if page_images:
        from app.knowledge.vision import describe_image

        for image_index, (page_index, image_bytes) in enumerate(page_images, start=1):
            try:
                description = describe_image(
                    image_bytes, f"{source} 第 {page_index} 页图片 {image_index}"
                )
                if description:
                    documents.append(
                        _document(
                            description,
                            source,
                            page=page_index,
                            section=f"第 {page_index} 页图片 {image_index}",
                            chunk_type="image_vision",
                        )
                    )
            except Exception as error:
                logger.warning("多模态图片解析失败，保留其他解析结果: %s", error)
    return documents


def _parse_docx(path: Path, source: str) -> list[Document]:
    from docx import Document as DocxDocument

    docx = DocxDocument(str(path))
    output: list[Document] = []
    section = ""
    for paragraph in docx.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = paragraph.style.name.lower() if paragraph.style else ""
        chunk_type = "title" if "heading" in style or "标题" in style else "text"
        if chunk_type == "title":
            section = text
        output.append(_document(text, source, section=section, chunk_type=chunk_type))
    for index, table in enumerate(docx.tables, start=1):
        rows = [
            " | ".join(cell.text.strip() for cell in row.cells) for row in table.rows
        ]
        output.append(
            _document(
                "\n".join(rows), source, section=f"表格 {index}", chunk_type="table"
            )
        )
    return output


def _parse_pptx(path: Path, source: str) -> list[Document]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    output: list[Document] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        title = (
            slide.shapes.title.text.strip()
            if slide.shapes.title
            else f"第 {slide_index} 页"
        )
        body: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                rows = [
                    " | ".join(cell.text.strip() for cell in row.cells)
                    for row in shape.table.rows
                ]
                output.append(
                    _document(
                        "\n".join(rows),
                        source,
                        page=slide_index,
                        section=title,
                        chunk_type="table",
                    )
                )
            elif (
                hasattr(shape, "text")
                and shape.text.strip()
                and shape is not slide.shapes.title
            ):
                body.append(shape.text.strip())
        output.append(
            _document(
                title, source, page=slide_index, section=title, chunk_type="title"
            )
        )
        if body:
            output.append(
                _document(
                    "\n".join(body),
                    source,
                    page=slide_index,
                    section=title,
                    chunk_type="text",
                )
            )
    return output


def _parse_workbook(path: Path, source: str) -> list[Document]:
    from openpyxl import load_workbook

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    output: list[Document] = []
    for sheet in workbook.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(values):
                rows.append(" | ".join(values))
        if rows:
            output.append(
                _document(
                    "\n".join(rows), source, section=sheet.title, chunk_type="table"
                )
            )
    return output


def _parse_csv(path: Path, source: str) -> list[Document]:
    with path.open("r", encoding="utf-8-sig", newline="") as file:
        rows = [" | ".join(row) for row in csv.reader(file)]
    return [_document("\n".join(rows), source, section=path.stem, chunk_type="table")]


def _parse_plain(path: Path, source: str) -> list[Document]:
    content = path.read_text(encoding="utf-8-sig", errors="replace")
    if path.suffix.lower() in {".html", ".htm"}:
        parser = _TextHTMLParser()
        parser.feed(content)
        content = "\n\n".join(parser.parts)
    return [_document(content, source, section=path.stem, chunk_type="text")]


def _parse_office_images(path: Path, source: str) -> list[Document]:
    if not settings.multimodal_enabled:
        return []
    from app.knowledge.vision import describe_image

    output: list[Document] = []
    try:
        with zipfile.ZipFile(path) as archive:
            media = [
                name
                for name in archive.namelist()
                if "/media/" in name
                and Path(name).suffix.lower() in {".png", ".jpg", ".jpeg", ".webp"}
            ]
            for index, name in enumerate(media[: settings.multimodal_max_images], start=1):
                image_bytes = archive.read(name)
                if len(image_bytes) < 10_000:
                    continue
                description = describe_image(image_bytes, f"{source} 内嵌图片 {index}")
                if description:
                    output.append(
                        _document(
                            description,
                            source,
                            section=f"内嵌图片 {index}",
                            chunk_type="image_vision",
                        )
                    )
    except (zipfile.BadZipFile, KeyError) as error:
        logger.warning("Office 内嵌图片读取失败: %s", error)
    return output


def parse_document(file_path: str, source_name: str | None = None) -> list[Document]:
    path = Path(file_path)
    extension = path.suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValueError(f"不支持的文件格式 {extension}，当前支持: {supported}")
    source = source_name or path.name
    parser = {
        ".pdf": _parse_pdf,
        ".docx": _parse_docx,
        ".pptx": _parse_pptx,
        ".xlsx": _parse_workbook,
        ".xlsm": _parse_workbook,
        ".csv": _parse_csv,
    }.get(extension, _parse_plain)
    documents = parser(path, source)
    if extension in {".docx", ".pptx", ".xlsx", ".xlsm"}:
        documents.extend(_parse_office_images(path, source))
    return documents


def ingest_document(
    user_id: str,
    kb_id: str,
    file_path: str,
    source_name: str | None = None,
    document_id: str | None = None,
) -> str:
    path = Path(file_path)
    source = source_name or path.name
    parsed = parse_document(file_path, source)
    chunks = chunk_documents(parsed)
    valid_chunks = [chunk for chunk in chunks if chunk.page_content.strip()]
    if not valid_chunks:
        raise ValueError("文档中没有可检索的文本；扫描件请启用 OCR 后重试")

    if document_id is None:
        digest = hashlib.sha256(path.read_bytes()).hexdigest()[:16]
        document_id = f"doc_{digest}"
    collection = vector_store.get_or_create_collection(user_id, kb_id)
    try:
        collection.delete(where={"document_id": document_id})
    except Exception:
        pass

    ids = [f"{document_id}_{index:05d}" for index in range(len(valid_chunks))]
    metadatas = []
    for chunk in valid_chunks:
        metadata = {
            "document_id": document_id,
            "source": source,
            "page": int(chunk.metadata.get("page", 0) or 0),
            "section": str(chunk.metadata.get("section", "")),
            "chunk_type": str(chunk.metadata.get("chunk_type", "text")),
            "parser": str(chunk.metadata.get("parser", "")),
        }
        metadatas.append(metadata)
    collection.upsert(
        ids=ids,
        documents=[chunk.page_content for chunk in valid_chunks],
        metadatas=metadatas,
    )
    logger.info(
        "文档入库完成: user=%s kb=%s document=%s source=%s chunks=%s",
        user_id,
        kb_id,
        document_id,
        source,
        len(valid_chunks),
    )
    return document_id
